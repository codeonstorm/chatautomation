import os
from pathlib import Path
from app.classes.document_splitter import DocumentSplitter
from app.classes.document_embedding import DocumentEmbedder
from app.classes.vector_db import VectorDB
from app.classes.file_helper import FileHelper
from uuid import UUID

from sqlmodel import select
from sqlmodel import Session
from app.core.database import get_session
from app.models.taskstatus import TaskStatus, TaskStageEnum

# from PyPDF2 import PdfReader
import docx
import openpyxl
import pptx
import mimetypes


class Ingestion:

    SUPPORTED_EXTENSIONS = ["docx", "txt", "xlsx", "ppt", "pptx"]

    def __init__(self, taskid: int, service_id: int, chatbot_uuid: UUID):
        self.taskid = taskid
        self.service_id = service_id
        self.chatbot_uuid = chatbot_uuid
        self.splitter = DocumentSplitter()
        self.embedder = DocumentEmbedder()
        self.vector_db = VectorDB(f"{str(chatbot_uuid)}")

    def read_file_content(self, file_path: Path, extension: str):
        """Extract text from various file types."""
        try:
            if extension == "pdf":
                pass
                # return self._read_pdf(file_path)
            elif extension == "docx":
                return self._read_docx(file_path)
            elif extension == "txt":
                return self._read_txt(file_path)
            elif extension == "xlsx":
                return self._read_excel(file_path)
            elif extension in ["ppt", "pptx"]:
                return self._read_ppt(file_path)
        except Exception as e:
            print(f"Error reading {file_path.name}: {e}")
            return ""

    # def _read_pdf(self, file_path):
    #     reader = PdfReader(file_path)
    #     return "\n".join([page.extract_text() or "" for page in reader.pages])

    def _read_docx(self, file_path):
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])

    def _read_txt(self, file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    def _read_excel(self, file_path):
        wb = openpyxl.load_workbook(file_path)
        text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text.append(" | ".join([str(cell) if cell else "" for cell in row]))
        return "\n".join(text)

    def _read_ppt(self, file_path):
        prs = pptx.Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def ingest(self):
        db: Session = next(get_session())
        upload_path = Path("uploads") / str(self.service_id)
        file_helper = FileHelper(upload_path)
        files = file_helper.get_file_details()

        total_target = len(files)
        target_processed = 0

        for file_info in files:
            # update progress...
            target_processed += 1
            progress = (target_processed / total_target) * 100

            statement = select(TaskStatus).where(TaskStatus.id == self.taskid)
            task = db.exec(statement)
            task = task.one()

            if progress == 100:
                task.status = TaskStageEnum.completed
            else:
                task.status = TaskStageEnum.in_progress
            task.progess = progress
            db.add(task)
            db.commit()

            extension = file_info["extension"].lower()
            if extension not in self.SUPPORTED_EXTENSIONS:
                print(f"Skipping unsupported file: {file_info['name']}")
                continue

            file_path = upload_path / file_info["name"]
            text = self.read_file_content(file_path, extension)

            if not text.strip():
                print(f"No text extracted from: {file_info['name']}")
                continue

            chunks = self.splitter.split_text(text)
            embeddings = self.embedder.embedding_model.encode(chunks)

            for chunk, embedding in zip(chunks, embeddings):
                self.vector_db.insert_vector(embedding.tolist(), chunk)

        self.vector_db.close()
